#!/usr/bin/env python3
"""Generate CTC Mobile Wishlist Business Case PowerPoint deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CT_RED = RGBColor(0xD5, 0x2B, 0x1E)
CT_DARK = RGBColor(0x33, 0x33, 0x33)
CT_GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
SUCCESS_GREEN = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=CT_DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, color=CT_DARK, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    return p

def slide_header(slide, title, subtitle=None):
    add_box(slide, Inches(0), Inches(0), W, Inches(1.1), CT_RED)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7), title, 32, WHITE, True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.4), subtitle, 16, RGBColor(0xFF,0xCC,0xCC))
    # Footer
    add_text(slide, Inches(0.6), Inches(7.0), Inches(5), Inches(0.4), "EPAM Systems | Confidential", 10, CT_GREY)
    add_text(slide, Inches(9), Inches(7.0), Inches(4), Inches(0.4), "CTC Mobile Wishlist — Business Case", 10, CT_GREY, alignment=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_box(s, Inches(0), Inches(0), W, Inches(3.2), CT_RED)
add_text(s, Inches(0.8), Inches(0.8), Inches(11), Inches(1.2), "CTC Mobile Wishlist", 48, WHITE, True)
add_text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.8), "Scan. Save. Share. — Turning In-Store Browsing Into Revenue", 24, RGBColor(0xFF,0xCC,0xCC))
add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5), "Hackathon POC  |  Business Case & Implementation Plan", 20, CT_DARK, True)
tf = add_text(s, Inches(0.8), Inches(4.4), Inches(11), Inches(2.5), "Kamal Syed", 20, CT_DARK, True)
add_para(tf, "Director of Program Management and Delivery, EPAM Systems", 14, CT_GREY)
add_para(tf, "", 8, CT_GREY)
add_para(tf, "Prepared for Canadian Tire Corporation  |  April 2026", 14, CT_GREY)
add_para(tf, "Agentic AI SDLC powered by EPAM EliteA", 14, CT_RED, True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Agenda")
items = [
    ("01", "The Opportunity — What's missing in CTC's mobile experience"),
    ("02", "The Concept — Scan, Save, Share wishlist feature"),
    ("03", "CTC Revenue Landscape — Current state by banner"),
    ("04", "Revenue Uplift Model — $50M–$548M incremental opportunity"),
    ("05", "Strategic Benefits — Beyond revenue: data, loyalty, differentiation"),
    ("06", "Competitive Landscape — First-mover advantage in Canada"),
    ("07", "Implementation Plan — Agentic AI SDLC with EPAM EliteA"),
    ("08", "Investment & ROI — Cost, timeline, and payback"),
    ("09", "POC Demo — What we built at the hackathon"),
    ("10", "Recommendation & Next Steps"),
]
for i, (num, text) in enumerate(items):
    y = Inches(1.4) + Inches(i * 0.55)
    add_text(s, Inches(1.0), y, Inches(0.6), Inches(0.45), num, 20, CT_RED, True)
    add_text(s, Inches(1.7), y, Inches(9), Inches(0.45), text, 18, CT_DARK)

print("Slides 1-2 done")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: The Opportunity
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "The Opportunity", "What's missing in CTC's mobile experience")

add_text(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5), "The Problem", 24, CT_RED, True)
gaps = [
    "No wishlist feature in the Canadian Tire mobile app",
    "No way to bridge in-store product discovery to digital intent",
    "No gift sharing / collaborative shopping functionality", 
    "16M Triangle members generate purchase data — but not intent data",
    "Cross-banner potential (CT + SportChek + Mark's) completely untapped",
]
tf = add_text(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.0), "", 16, CT_DARK)
for g in gaps:
    add_para(tf, f"    {g}", 16, CT_DARK, space_before=Pt(10))

add_text(s, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5), "The Insight", 24, SUCCESS_GREEN, True)
insights = [
    "Wishlist users convert 25-40% better than non-users",
    "Wishlist users spend 15-25% more per transaction",
    "Price-drop alerts on wishlisted items get 5-8x higher CTR",
    "Gift wishlists see 60-70% fulfilment rates",
    "Wishlists recover 10-15% of would-be cart abandonments",
]
tf2 = add_text(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.0), "", 16, CT_DARK)
for ins in insights:
    add_para(tf2, f"    {ins}", 16, CT_DARK, space_before=Pt(10))

add_box(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.7), LIGHT_GREY)
add_text(s, Inches(1.0), Inches(6.1), Inches(11), Inches(0.5),
    "Canadian Tire has the infrastructure (16M loyalty members, 500+ stores, 10M+ app installs) — it just needs the feature.",
    14, CT_DARK, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: The Concept
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "The Concept", "Scan. Save. Share.")

features = [
    ("SCAN", "Scan any barcode on\na store shelf with your\nphone camera to instantly\nidentify the product", "In-Store\nDiscovery"),
    ("SAVE", "Add products to named\nwishlists from browsing\nthe catalog or scanning\nbarcodes", "Wishlist\nManagement"),
    ("SHARE", "Share your wishlist with\nfamily & friends from\nyour phone contacts for\ngift coordination", "Social\nCommerce"),
    ("FULFILL", "Recipients browse your\nlist, claim items they'll\nbuy — no duplicate\ngifts", "Gift\nFulfillment"),
]

for i, (title, desc, label) in enumerate(features):
    x = Inches(0.6 + i * 3.15)
    # Card background
    add_box(s, x, Inches(1.5), Inches(2.9), Inches(4.8), LIGHT_GREY)
    # Number circle
    circle = add_box(s, x + Inches(1.05), Inches(1.7), Inches(0.8), Inches(0.8), CT_RED)
    add_text(s, x + Inches(1.05), Inches(1.75), Inches(0.8), Inches(0.8), str(i+1), 28, WHITE, True, PP_ALIGN.CENTER)
    # Title
    add_text(s, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.5), title, 22, CT_RED, True, PP_ALIGN.CENTER)
    # Description
    add_text(s, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(1.5), desc, 14, CT_GREY, alignment=PP_ALIGN.CENTER)
    # Bottom label
    add_box(s, x + Inches(0.4), Inches(5.2), Inches(2.1), Inches(0.8), CT_RED)
    add_text(s, x + Inches(0.4), Inches(5.25), Inches(2.1), Inches(0.8), label, 13, WHITE, True, PP_ALIGN.CENTER)

print("Slides 3-4 done")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: CTC Revenue Landscape
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "CTC Revenue Landscape", "Sources: CTC FY 2024 Annual Report, Q4 2024 Earnings Release, CTC 2024 AIF (SEDAR+)")

# Left: Revenue by banner table
add_text(s, Inches(0.8), Inches(1.3), Inches(6), Inches(0.5), "Revenue by Banner", 20, CT_DARK, True)

headers = ["Banner", "Revenue", "Stores", "Avg. Transaction"]
col_x = [Inches(0.8), Inches(4.0), Inches(5.8), Inches(7.0)]
col_w = [Inches(3.2), Inches(1.8), Inches(1.2), Inches(2.0)]

# Header row
add_box(s, Inches(0.8), Inches(1.8), Inches(8.2), Inches(0.45), CT_RED)
for hdr, cx, cw in zip(headers, col_x, col_w):
    add_text(s, cx, Inches(1.82), cw, Inches(0.4), hdr, 13, WHITE, True)

rows = [
    ("Canadian Tire Retail", "$9.8B", "500+", "$45-65"),
    ("SportChek / Sports Experts", "$2.2B", "185", "$70-90"),
    ("Mark's / L'Equipeur", "$1.4B", "385", "$55-75"),
    ("Helly Hansen", "$1.0B", "65+", "N/A"),
    ("Party City Canada", "$175M", "69+", "$30-50"),
    ("PartSource", "$76M", "~95", "$50-80"),
    ("CT Financial Services (CT Bank)", "$1.9B", "—", "—"),
]
# Source footnote
add_text(s, Inches(0.8), Inches(5.2), Inches(8), Inches(0.6),
    "Sources: CTC Q4 2024 Earnings Release (Feb 2025); CTC 2024 Annual Information Form (SEDAR+);\n"
    "Party City: CTC acquisition ($174.4M, Aug 2019, CBC); PartSource: ZoomInfo/Growjo est.; Store counts: CTC 2024 AIF.",
    8, CT_GREY)
for ri, (b, r, st, atv) in enumerate(rows):
    y = Inches(2.3 + ri * 0.36)
    bg_c = LIGHT_GREY if ri % 2 == 0 else WHITE
    add_box(s, Inches(0.8), y, Inches(8.2), Inches(0.42), bg_c)
    vals = [b, r, st, atv]
    for v, cx, cw in zip(vals, col_x, col_w):
        add_text(s, cx, y, cw, Inches(0.4), v, 12, CT_DARK)

# Right: Key metrics
add_text(s, Inches(9.5), Inches(1.3), Inches(3.5), Inches(0.5), "Digital & Loyalty", 20, CT_DARK, True)
metrics = [
    ("16M", "Triangle Rewards\nmembers"),
    ("65-70%", "Transaction\npenetration"),
    ("10M+", "App downloads\n(cumulative)"),
    ("$2.2B", "E-commerce\nrevenue"),
    ("14-16%", "Digital share\nof retail"),
]
for i, (val, label) in enumerate(metrics):
    y = Inches(1.9 + i * 1.0)
    add_text(s, Inches(9.5), y, Inches(1.3), Inches(0.5), val, 24, CT_RED, True)
    add_text(s, Inches(10.8), y, Inches(2.2), Inches(0.8), label, 12, CT_GREY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: Revenue Uplift Model
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Revenue Uplift Model", "Four independent revenue channels — $50M to $548M annually")

channels = [
    ("Direct Wishlist Revenue", "Higher conversion + AOV\nfrom wishlist users", "$21.6M", "$96.0M", "$315.0M"),
    ("Gift Sharing / Fulfillment", "Shared wishlists drive\nfull-price gift purchases", "$10.8M", "$45.0M", "$140.0M"),
    ("Abandonment Recovery", "Wishlists catch items that\nwould otherwise be lost", "$8.3M", "$21.0M", "$39.0M"),
    ("Re-engagement Alerts", "Price-drop / back-in-stock\non wishlisted items", "$10.0M", "$27.5M", "$54.0M"),
]

# Column headers
add_box(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(0.5), CT_RED)
ch = [("Revenue Channel", 0.8, 3.0), ("Driver", 3.8, 2.5), ("Conservative", 6.5, 1.8), ("Moderate", 8.5, 1.8), ("Aggressive", 10.5, 1.8)]
for label, x, w in ch:
    add_text(s, Inches(x), Inches(1.32), Inches(w), Inches(0.45), label, 14, WHITE, True)

for ri, (name, driver, con, mod, agg) in enumerate(channels):
    y = Inches(1.9 + ri * 1.05)
    bg_c = LIGHT_GREY if ri % 2 == 0 else WHITE
    add_box(s, Inches(0.6), y, Inches(12.0), Inches(0.95), bg_c)
    add_text(s, Inches(0.8), y + Inches(0.1), Inches(3.0), Inches(0.8), name, 15, CT_DARK, True)
    add_text(s, Inches(3.8), y + Inches(0.1), Inches(2.5), Inches(0.8), driver, 11, CT_GREY)
    add_text(s, Inches(6.5), y + Inches(0.2), Inches(1.8), Inches(0.5), con, 16, CT_DARK, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(8.5), y + Inches(0.2), Inches(1.8), Inches(0.5), mod, 16, CT_RED, True, PP_ALIGN.CENTER)
    add_text(s, Inches(10.5), y + Inches(0.2), Inches(1.8), Inches(0.5), agg, 16, CT_DARK, alignment=PP_ALIGN.CENTER)

# Total row
ty = Inches(6.1)
add_box(s, Inches(0.6), ty, Inches(12.0), Inches(0.7), CT_DARK)
add_text(s, Inches(0.8), ty + Inches(0.1), Inches(3.0), Inches(0.5), "TOTAL", 16, WHITE, True)
add_text(s, Inches(6.5), ty + Inches(0.1), Inches(1.8), Inches(0.5), "$50.7M", 16, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(8.5), ty + Inches(0.1), Inches(1.8), Inches(0.5), "$189.5M", 18, RGBColor(0xFF,0xCC,0xCC), True, PP_ALIGN.CENTER)
add_text(s, Inches(10.5), ty + Inches(0.1), Inches(1.8), Inches(0.5), "$548.0M", 16, WHITE, True, PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.4),
    "All amounts in CAD. Moderate scenario = ~1.3% lift on $14.5B retail revenue. Assumptions: 25% adoption, +$12 AOV, 8 txn/year.",
    12, CT_GREY, alignment=PP_ALIGN.LEFT)

print("Slides 5-6 done")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: Revenue by Banner
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Revenue Uplift by Banner", "Moderate scenario — $189.5M across all CTC banners  |  Revenue data: CTC FY2024 Annual Report & AIF")

banners = [
    ("Canadian Tire", "$9.8B", "~15%", "$112M", "Largest user base, broadest catalog,\nstrongest barcode scan use case", 7.5),
    ("SportChek", "$2.2B", "~18%", "$36M", "Higher AOV, strong gift-giving\ncategory (sports equipment)", 4.0),
    ("Mark's", "$1.4B", "~12%", "$21M", "Workwear wishlists, seasonal\napparel gifting", 2.4),
    ("Helly Hansen", "$1.0B", "~25%", "$14M", "Premium outdoor gear, high wishlist\naffinity, global DTC", 1.7),
    ("Party City", "$175M", "~8%", "$4M", "Celebration wishlists — birthday &\nholiday party planning", 0.6),
    ("PartSource", "$76M", "~5%", "$2.5M", "Auto parts wishlists — save parts\nfor scheduled maintenance", 0.4),
]

for i, (name, rev, ecom, uplift, notes, bar_w) in enumerate(banners):
    y = Inches(1.4 + i * 0.9)
    # Banner name and stats
    add_text(s, Inches(0.8), y, Inches(2.5), Inches(0.4), name, 18, CT_DARK, True)
    add_text(s, Inches(0.8), y + Inches(0.4), Inches(2.5), Inches(0.3), f"Revenue: {rev}  |  E-comm: {ecom}", 11, CT_GREY)
    # Bar
    add_box(s, Inches(3.5), y + Inches(0.1), Inches(bar_w * 0.85), Inches(0.5), CT_RED)
    add_text(s, Inches(3.5) + Inches(bar_w * 0.85) + Inches(0.15), y + Inches(0.1), Inches(1.5), Inches(0.5), uplift, 18, CT_RED, True)
    # Notes
    add_text(s, Inches(3.5), y + Inches(0.65), Inches(8), Inches(0.6), notes, 11, CT_GREY)

# Total
add_box(s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.5), LIGHT_GREY)
add_text(s, Inches(0.8), Inches(6.42), Inches(4), Inches(0.45), "Total Moderate Uplift:", 16, CT_DARK, True)
add_text(s, Inches(5.0), Inches(6.42), Inches(3), Inches(0.45), "$189.5M / year", 16, CT_RED, True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: Strategic Benefits
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Strategic Benefits Beyond Revenue")

add_text(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "Consumer Benefits", 20, CT_RED, True)
consumer = [
    ("Save for Later", "Bookmark items without cart commitment"),
    ("In-Store to Digital Bridge", "Scan shelf products, buy later anywhere"),
    ("Gift Coordination", "Share lists, eliminate duplicate gifts"),
    ("Price Monitoring", "Alerts when wishlisted items go on sale"),
    ("Cross-Device Continuity", "Start on phone in-store, buy on desktop"),
]
for i, (title, desc) in enumerate(consumer):
    y = Inches(1.9 + i * 0.75)
    add_text(s, Inches(0.8), y, Inches(5.5), Inches(0.35), title, 14, CT_DARK, True)
    add_text(s, Inches(0.8), y + Inches(0.3), Inches(5.5), Inches(0.35), desc, 12, CT_GREY)

add_text(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5), "CTC Business Benefits", 20, SUCCESS_GREEN, True)
business = [
    ("First-Party Intent Data", "What customers WANT — not just what they buy"),
    ("Triangle Enrichment", "Richest customer intent dataset in Canadian retail"),
    ("Seasonal Revenue Capture", "Convert holiday browsing into purchase pipelines"),
    ("Dealer Inventory Signals", "Aggregated wishlist data improves local stocking"),
    ("Cross-Banner Discovery", "Unified wishlists across all 6 CTC banners"),
    ("Competitive Moat", "Only Canadian retailer with scan-to-wishlist + sharing"),
]
for i, (title, desc) in enumerate(business):
    y = Inches(1.9 + i * 0.75)
    add_text(s, Inches(7.0), y, Inches(5.5), Inches(0.35), title, 14, CT_DARK, True)
    add_text(s, Inches(7.0), y + Inches(0.3), Inches(5.5), Inches(0.35), desc, 12, CT_GREY)

print("Slides 7-8 done")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: Competitive Landscape
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Competitive Landscape", "CTC would be the only Canadian retailer with all four capabilities")

# Table
headers2 = ["Retailer", "Wishlist", "Barcode Scan", "Gift Sharing", "Cross-Banner"]
cx2 = [Inches(0.8), Inches(4.5), Inches(6.3), Inches(8.1), Inches(10.0)]
cw2 = [Inches(3.7), Inches(1.8), Inches(1.8), Inches(1.9), Inches(2.0)]

add_box(s, Inches(0.6), Inches(1.5), Inches(11.8), Inches(0.5), CT_DARK)
for h, cx, cw in zip(headers2, cx2, cw2):
    add_text(s, cx, Inches(1.52), cw, Inches(0.45), h, 14, WHITE, True)

comp_rows = [
    ("Canadian Tire (proposed)", "Yes", "Yes", "Yes", "Yes", True),
    ("Amazon.ca", "Yes", "No", "Yes", "N/A", False),
    ("Walmart Canada", "Basic", "No", "No", "No", False),
    ("Costco Canada", "No", "No", "No", "N/A", False),
    ("Home Depot Canada", "Basic", "No", "No", "N/A", False),
    ("Best Buy Canada", "Yes", "No", "Limited", "N/A", False),
]
for ri, (name, w, b, g, cb, highlight) in enumerate(comp_rows):
    y = Inches(2.1 + ri * 0.6)
    bg_c = RGBColor(0xFF, 0xEB, 0xEE) if highlight else (LIGHT_GREY if ri % 2 == 0 else WHITE)
    add_box(s, Inches(0.6), y, Inches(11.8), Inches(0.55), bg_c)
    n_color = CT_RED if highlight else CT_DARK
    add_text(s, Inches(0.8), y + Inches(0.05), Inches(3.7), Inches(0.45), name, 14, n_color, highlight)
    for val, cx, cw in zip([w, b, g, cb], cx2[1:], cw2[1:]):
        v_color = SUCCESS_GREEN if val == "Yes" else (CT_RED if val == "No" else CT_GREY)
        add_text(s, cx, y + Inches(0.05), cw, Inches(0.45), val, 14, v_color, val=="Yes", PP_ALIGN.CENTER)

add_box(s, Inches(0.8), Inches(5.9), Inches(11.0), Inches(0.7), LIGHT_GREY)
add_text(s, Inches(1.0), Inches(6.0), Inches(10.5), Inches(0.5),
    "No Canadian competitor combines all four capabilities. CTC has a clear first-mover advantage.",
    15, CT_DARK, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: Agentic AI SDLC with EPAM EliteA
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Implementation Approach", "Fully Agentic AI SDLC powered by EPAM EliteA")

add_text(s, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
    "EPAM EliteA: AI-Augmented Software Development at Scale", 20, CT_DARK, True)

tf = add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
    "EliteA is EPAM's proprietary agentic AI platform that orchestrates the full SDLC — from requirements to deployment — "
    "using specialized AI agents supervised by senior engineers. This approach delivers 3-5x velocity improvement over "
    "traditional development while maintaining enterprise-grade quality.", 14, CT_GREY)

phases = [
    ("Requirements\n& Architecture", "AI agents analyze\nbusiness requirements,\ngenerate architecture\ndocs, data schemas,\nand API contracts", "1-2 weeks"),
    ("Agentic\nDevelopment", "Coding agents generate\nfeatures in parallel.\nHuman engineers review,\nguide, and approve\nall output", "3-4 weeks"),
    ("AI-Assisted\nQA & Testing", "Test agents generate\nunit/integration/E2E\ntests. Coverage agents\nenforce 80%+ coverage\nthresholds", "1-2 weeks"),
    ("Integration\n& Hardening", "Integration with CTC\nbackend APIs, Triangle\nRewards, product catalog.\nSecurity & perf testing", "2-3 weeks"),
    ("Deployment\n& Launch", "CI/CD pipeline, staged\nrollout to 50 pilot\nstores, monitoring,\nA/B testing framework", "1-2 weeks"),
]

for i, (title, desc, duration) in enumerate(phases):
    x = Inches(0.4 + i * 2.5)
    add_box(s, x, Inches(3.0), Inches(2.3), Inches(3.5), LIGHT_GREY)
    # Phase number
    add_box(s, x + Inches(0.8), Inches(3.1), Inches(0.7), Inches(0.7), CT_RED)
    add_text(s, x + Inches(0.8), Inches(3.15), Inches(0.7), Inches(0.7), str(i+1), 22, WHITE, True, PP_ALIGN.CENTER)
    # Title
    add_text(s, x + Inches(0.1), Inches(3.9), Inches(2.1), Inches(0.7), title, 13, CT_DARK, True, PP_ALIGN.CENTER)
    # Description
    add_text(s, x + Inches(0.1), Inches(4.6), Inches(2.1), Inches(1.5), desc, 10, CT_GREY, alignment=PP_ALIGN.CENTER)
    # Duration badge
    add_box(s, x + Inches(0.5), Inches(6.1), Inches(1.3), Inches(0.35), CT_DARK)
    add_text(s, x + Inches(0.5), Inches(6.1), Inches(1.3), Inches(0.35), duration, 11, WHITE, True, PP_ALIGN.CENTER)

print("Slides 9-10 done")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: Resource Plan Comparison
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Resource Plan Comparison", "All amounts in CAD  |  Onshore: DM $200/h, BA $160/h, SA $190/h  |  160h/month")

# ── LEFT: Traditional ──
add_box(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.5), CT_DARK)
add_text(s, Inches(0.7), Inches(1.32), Inches(5.5), Inches(0.45), "Traditional Development Model", 16, WHITE, True)

# Header row
add_box(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(0.32), LIGHT_GREY)
for label, lx, lw in [("Role", 0.6, 2.0), ("HC", 2.65, 0.35), ("Shore", 3.05, 0.55), ("Rate/h", 3.65, 0.6), ("Duration", 4.35, 0.8), ("Cost CAD", 5.2, 1.2)]:
    add_text(s, Inches(lx), Inches(1.85), Inches(lw), Inches(0.3), label, 8, CT_GREY, True)

trad_roles = [
    ("Delivery Manager",     "1", "On", "$200", "6-9 mo",  "$192K-288K"),
    ("Business Analyst",     "1", "On", "$160", "2-3 mo",  "$51K-77K"),
    ("Solution Architect",   "1", "On", "$190", "3-4 mo",  "$91K-122K"),
    ("Sr Backend Developer", "2", "Off", "$130", "5-6 mo", "$208K-250K"),
    ("Sr Mobile Developer",  "2", "Off", "$130", "5-6 mo", "$208K-250K"),
    ("Functional Tester",    "1", "Off", "$95",  "3-4 mo", "$46K-61K"),
    ("Automation Tester",    "1", "Off", "$110", "3-4 mo", "$53K-70K"),
    ("DevOps / SRE",         "1", "Off", "$140", "2.5-3 mo", "$56K-67K"),
]

for ri, (role, hc, shore, rate, dur, cost) in enumerate(trad_roles):
    y = Inches(2.2 + ri * 0.3)
    bg_c = LIGHT_GREY if ri % 2 == 0 else WHITE
    add_box(s, Inches(0.5), y, Inches(6.0), Inches(0.3), bg_c)
    shore_c = CT_RED if shore == "On" else RGBColor(0x15, 0x65, 0xC0)
    for val, vx, vw, vc in zip(
        [role, hc, shore, rate, dur, cost],
        [0.6, 2.65, 3.05, 3.65, 4.35, 5.2],
        [2.0, 0.35, 0.55, 0.6, 0.8, 1.2],
        [CT_DARK, CT_DARK, shore_c, CT_DARK, CT_DARK, CT_DARK]):
        add_text(s, Inches(vx), y, Inches(vw), Inches(0.28), val, 9, vc, val == shore and shore == "On")

# Totals
add_box(s, Inches(0.5), Inches(4.65), Inches(6.0), Inches(0.4), CT_DARK)
add_text(s, Inches(0.6), Inches(4.67), Inches(1.8), Inches(0.35), "Total: 10 FTEs", 11, WHITE, True)
add_text(s, Inches(2.7), Inches(4.67), Inches(1.5), Inches(0.35), "3 onshore + 7 offshore", 9, RGBColor(0xCC,0xCC,0xCC))
add_text(s, Inches(4.35), Inches(4.67), Inches(0.8), Inches(0.35), "6-9 months", 11, WHITE, True)
add_text(s, Inches(5.2), Inches(4.67), Inches(1.2), Inches(0.35), "$905K-1.18M", 11, WHITE, True)

# Timeline visual
add_text(s, Inches(0.6), Inches(5.2), Inches(5.5), Inches(0.25), "Sequential — each phase waits for the previous:", 9, CT_GREY)
trad_phases = [("Arch", 1.0), ("Backend", 2.0), ("Mobile", 2.0), ("QA", 1.1), ("Launch", 0.65)]
px = Inches(0.6)
for i, (label, w) in enumerate(trad_phases):
    add_box(s, px, Inches(5.5), Inches(w), Inches(0.3), [CT_RED, RGBColor(0x15,0x65,0xC0), RGBColor(0xF5,0x7C,0x00), SUCCESS_GREEN, CT_DARK][i])
    add_text(s, px, Inches(5.5), Inches(w), Inches(0.3), label, 8, WHITE, True, PP_ALIGN.CENTER)
    px = px + Inches(w + 0.04)

# ── RIGHT: EliteA ──
add_box(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.5), CT_RED)
add_text(s, Inches(7.0), Inches(1.32), Inches(5.5), Inches(0.45), "EPAM EliteA Agentic AI Model", 16, WHITE, True)

# Header row
add_box(s, Inches(6.8), Inches(1.85), Inches(6.0), Inches(0.32), LIGHT_GREY)
for label, lx, lw in [("Role", 6.9, 2.0), ("HC", 8.95, 0.8), ("Shore", 9.8, 0.5), ("Rate/h", 10.35, 0.55), ("Duration", 10.95, 0.75), ("Cost CAD", 11.75, 1.0)]:
    add_text(s, Inches(lx), Inches(1.85), Inches(lw), Inches(0.3), label, 8, CT_GREY, True)

elite_roles = [
    ("Delivery Manager",       "1",          "On",  "$200",  "15 wks",  "$120K"),
    ("Business Analyst",       "1 (PT)",     "On",  "$160",  "4 wks",   "$26K"),
    ("Solution Architect",     "1",          "On",  "$190",  "6 wks",   "$46K"),
    ("Sr Full-Stack Engineer", "1",          "Off", "$130",  "12 wks",  "$62K"),
    ("Sr Mobile Engineer",     "1",          "Off", "$130",  "10 wks",  "$52K"),
    ("AI Agent Cluster (Dev)", "5-8 agents", "AI",  "—",     "4-6 wks", "$45K"),
    ("AI Agent Cluster (QA)",  "2-3 agents", "AI",  "—",     "2-3 wks", "$20K"),
    ("Human QA Reviewer",      "1 (PT)",     "Off", "$100",  "5 wks",   "$20K"),
    ("DevOps (human + AI)",    "0.5 + AI",   "Off", "$140",  "3 wks",   "$17K"),
]

for ri, (role, hc, shore, rate, dur, cost) in enumerate(elite_roles):
    y = Inches(2.2 + ri * 0.27)
    bg_c = LIGHT_GREY if ri % 2 == 0 else WHITE
    add_box(s, Inches(6.8), y, Inches(6.0), Inches(0.27), bg_c)
    if shore == "AI":
        shore_c = CT_RED
    elif shore == "On":
        shore_c = CT_RED
    else:
        shore_c = RGBColor(0x15, 0x65, 0xC0)
    is_ai = shore == "AI"
    for val, vx, vw, vc in zip(
        [role, hc, shore, rate, dur, cost],
        [6.9, 8.95, 9.8, 10.35, 10.95, 11.75],
        [2.0, 0.8, 0.5, 0.55, 0.75, 1.0],
        [CT_RED if is_ai else CT_DARK, CT_RED if is_ai else CT_DARK, shore_c, CT_DARK, CT_DARK, CT_DARK]):
        add_text(s, Inches(vx), y, Inches(vw), Inches(0.26), val, 8, vc, is_ai)

# Totals
add_box(s, Inches(6.8), Inches(4.65), Inches(6.0), Inches(0.4), CT_RED)
add_text(s, Inches(6.9), Inches(4.67), Inches(1.8), Inches(0.35), "5 humans + AI", 11, WHITE, True)
add_text(s, Inches(8.7), Inches(4.67), Inches(1.5), Inches(0.35), "3 onshore + 2 offshore", 9, RGBColor(0xFF,0xCC,0xCC))
add_text(s, Inches(10.95), Inches(4.67), Inches(0.75), Inches(0.35), "8-15 wks", 11, WHITE, True)
add_text(s, Inches(11.75), Inches(4.67), Inches(1.0), Inches(0.35), "$408K", 11, WHITE, True)

# Timeline visual
add_text(s, Inches(6.9), Inches(5.2), Inches(5.5), Inches(0.25), "Parallel — AI agents run concurrently with human supervision:", 9, CT_GREY)
elite_phases = [("Arch", 0.65), ("Backend + Mobile + QA (parallel)", 3.0), ("Integrate", 0.9), ("Launch", 0.55)]
px = Inches(6.9)
for i, (label, w) in enumerate(elite_phases):
    add_box(s, px, Inches(5.5), Inches(w), Inches(0.3), [CT_RED, RGBColor(0x15,0x65,0xC0), RGBColor(0xF5,0x7C,0x00), CT_DARK][i])
    add_text(s, px, Inches(5.5), Inches(w), Inches(0.3), label, 8, WHITE, True, PP_ALIGN.CENTER)
    px = px + Inches(w + 0.04)

# Bottom comparison callout
add_box(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.35), LIGHT_GREY)
add_text(s, Inches(0.7), Inches(6.0), Inches(11.8), Inches(0.33),
    "EliteA saves $500K-770K CAD (55-65%)  |  50% fewer human FTEs (5 vs 10)  |  4-6 months faster  |  AI handles ~70% of code, tests, and docs",
    11, CT_RED, True, PP_ALIGN.CENTER)

# Role legend
add_text(s, Inches(0.7), Inches(6.45), Inches(11.8), Inches(0.35),
    "On = Onshore (Canada)  |  Off = Offshore/Nearshore  |  AI = EPAM EliteA agent compute  |  PT = Part-time  |  All rates in CAD/hour",
    9, CT_GREY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: Investment & ROI
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Investment & ROI", "All amounts in CAD  |  Agentic AI SDLC reduces cost by ~60% and timeline by ~50%")

# Left: Cost comparison
add_text(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "Implementation Cost Comparison (CAD)", 18, CT_DARK, True)

# Traditional
add_text(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.35), "Traditional SDLC", 14, CT_GREY, True)
trad_items = [
    ("Design & Architecture", "$142K-199K", "4-6 weeks"),
    ("Backend Development", "$264K-317K", "8-12 weeks"),
    ("Mobile Development", "$264K-317K", "8-12 weeks"),
    ("QA (Functional + Automation)", "$99K-131K", "4-6 weeks"),
    ("DevOps / Launch", "$56K-67K", "2-4 weeks"),
    ("DM + BA Oversight", "$243K-365K", "Full duration"),
]
for i, (phase, cost, dur) in enumerate(trad_items):
    y = Inches(2.3 + i * 0.35)
    add_text(s, Inches(0.8), y, Inches(2.5), Inches(0.32), phase, 11, CT_DARK)
    add_text(s, Inches(3.3), y, Inches(1.2), Inches(0.32), cost, 11, CT_GREY, alignment=PP_ALIGN.RIGHT)
    add_text(s, Inches(4.6), y, Inches(1.2), Inches(0.32), dur, 11, CT_GREY)

add_box(s, Inches(0.8), Inches(4.45), Inches(5.0), Inches(0.4), LIGHT_GREY)
add_text(s, Inches(0.8), Inches(4.47), Inches(2.5), Inches(0.35), "Traditional Total:", 13, CT_DARK, True)
add_text(s, Inches(3.3), Inches(4.47), Inches(1.2), Inches(0.35), "$905K-1.18M", 13, CT_DARK, True, PP_ALIGN.RIGHT)
add_text(s, Inches(4.6), Inches(4.47), Inches(1.2), Inches(0.35), "6-9 months", 13, CT_DARK, True)

# EliteA
add_text(s, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.35), "EPAM EliteA Agentic AI SDLC", 14, CT_RED, True)
elite_items = [
    ("Onshore (DM + BA + SA)", "$192K", "Full / partial"),
    ("Offshore Engineers (2)", "$114K", "10-12 weeks"),
    ("AI Agent Compute", "$65K", "4-6 weeks"),
    ("Human QA + DevOps", "$37K", "3-5 weeks"),
]
for i, (phase, cost, dur) in enumerate(elite_items):
    y = Inches(5.5 + i * 0.32)
    add_text(s, Inches(0.8), y, Inches(2.5), Inches(0.3), phase, 11, CT_DARK)
    add_text(s, Inches(3.3), y, Inches(1.2), Inches(0.3), cost, 11, SUCCESS_GREEN, alignment=PP_ALIGN.RIGHT)
    add_text(s, Inches(4.6), y, Inches(1.2), Inches(0.3), dur, 11, SUCCESS_GREEN)

add_box(s, Inches(0.8), Inches(6.8), Inches(5.0), Inches(0.4), CT_RED)
add_text(s, Inches(0.8), Inches(6.82), Inches(2.5), Inches(0.35), "EliteA Total:", 13, WHITE, True)
add_text(s, Inches(3.3), Inches(6.82), Inches(1.2), Inches(0.35), "$408K", 13, WHITE, True, PP_ALIGN.RIGHT)
add_text(s, Inches(4.6), Inches(6.82), Inches(1.2), Inches(0.35), "8-15 weeks", 13, WHITE, True)

# Right: ROI metrics
add_text(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5), "ROI Analysis — Moderate Scenario (CAD)", 18, CT_DARK, True)

roi_metrics = [
    ("Year 1 Incremental Revenue", "$189.5M CAD"),
    ("Blended Gross Margin", "~35%"),
    ("Year 1 Gross Profit", "$66.3M CAD"),
    ("EliteA Build Investment", "$408K CAD"),
    ("Payback Period", "< 1 day"),
    ("Year 1 ROI", "16,150%"),
]
for i, (label, val) in enumerate(roi_metrics):
    y = Inches(2.0 + i * 0.7)
    add_text(s, Inches(7.0), y, Inches(3.5), Inches(0.35), label, 14, CT_DARK)
    is_highlight = i >= 4
    add_text(s, Inches(10.5), y, Inches(2.3), Inches(0.4), val, 20 if is_highlight else 16,
        CT_RED if is_highlight else CT_DARK, is_highlight, PP_ALIGN.RIGHT)

# Savings callout
add_box(s, Inches(7.0), Inches(6.2), Inches(5.5), Inches(0.7), LIGHT_GREY)
tf = add_text(s, Inches(7.2), Inches(6.25), Inches(5.1), Inches(0.6),
    "EliteA saves $500K-770K CAD and 4+ months vs. traditional", 14, CT_RED, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: Assumptions
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Key Assumptions", "All amounts in CAD  |  EliteA agentic AI SDLC cost and timeline estimates")

add_text(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "EliteA Delivery Assumptions", 18, CT_RED, True)

assumptions_left = [
    "3-5x developer velocity via AI pair programming and code generation",
    "Senior EPAM engineers supervise all AI-generated output (human-in-the-loop)",
    "AI agents handle ~70% of boilerplate code, tests, and documentation",
    "Human engineers focus on architecture, integration, edge cases, and review",
    "Onshore rates (CAD/h): DM $200, BA $160, SA $190",
    "Offshore rates (CAD/h): Sr Dev $130, QA $95-110, DevOps $140",
    "Team: 3 onshore (DM, BA, SA) + 2 offshore engineers + AI agent cluster",
    "Parallel workstreams: backend + mobile + QA run concurrently",
    "CTC provides API access, design assets, and Triangle sandbox within 2 weeks",
]

tf = add_text(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(5.0), "", 12, CT_DARK)
for a in assumptions_left:
    add_para(tf, f"  {a}", 12, CT_DARK, space_before=Pt(8))

add_text(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5), "Revenue Model Assumptions", 18, CT_RED, True)

assumptions_right = [
    "Mobile MAU baseline: 4M across all CTC banners",
    "Wishlist adoption: 25% of MAU (1M users) within Year 1",
    "Incremental AOV lift: +$12 CAD per transaction for wishlist users",
    "Average 8 transactions/year for engaged wishlist users",
    "Gift sharing: 30% of wishlist users share at least one list",
    "Average fulfilled gift value: $150 CAD per shared wishlist",
    "Abandonment recovery: 350K transactions saved at $60 avg",
    "Re-engagement CTR: 5-8x higher than generic promotional pushes",
    "Blended gross margin: 35% (mix of owned brands + national brands)",
    "No cannibalization assumed — wishlist drives net-new purchases",
]

tf2 = add_text(s, Inches(7.0), Inches(1.9), Inches(5.5), Inches(5.0), "", 12, CT_DARK)
for a in assumptions_right:
    add_para(tf2, f"  {a}", 12, CT_DARK, space_before=Pt(8))

print("Slides 11-12 done")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13: POC Demo
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "POC Demonstration", "Built in <48 hours at the EPAM-CTC Hackathon using agentic AI")

add_text(s, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
    "What We Built", 20, CT_DARK, True)
tf = add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "A fully functional React Native + Expo mobile app demonstrating the complete wishlist journey — built entirely with EPAM EliteA's agentic AI workflow.", 14, CT_GREY)

screens = [
    ("Home", "Personalized landing\nwith quick access to\nwishlists and recent\nscanned items"),
    ("Catalog", "Browsable product grid\nwith category filters\nand real-time search"),
    ("Barcode\nScanner", "Camera-based barcode\nscanning with instant\nproduct identification"),
    ("Wishlist\nManager", "Create, name, and\nmanage multiple\nwishlists with items"),
    ("Share &\nFulfill", "Share with contacts,\nrecipients claim items\nthey'll purchase"),
]

for i, (title, desc) in enumerate(screens):
    x = Inches(0.5 + i * 2.5)
    # Phone mockup shape
    add_box(s, x, Inches(2.6), Inches(2.2), Inches(3.8), LIGHT_GREY)
    add_box(s, x + Inches(0.1), Inches(2.7), Inches(2.0), Inches(0.4), CT_RED)
    add_text(s, x + Inches(0.1), Inches(2.72), Inches(2.0), Inches(0.4), title, 13, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.3), Inches(1.9), Inches(2.5), desc, 11, CT_GREY, alignment=PP_ALIGN.CENTER)

add_box(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4), LIGHT_GREY)
add_text(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.4),
    "Tech: React Native + Expo  |  TypeScript  |  Local mock data  |  iOS + Android  |  Built with EPAM EliteA agentic AI",
    12, CT_GREY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14: KPIs
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Success Metrics", "Year 1 KPIs for the fully realized feature")

kpis = [
    ("Wishlist Adoption", ">=20%", "of mobile MAU", "Core adoption target"),
    ("Avg Items/Wishlist", ">=5", "items", "Engagement depth"),
    ("Wishlist-to-Purchase", ">=15%", "conversion rate", "Revenue validation"),
    ("Gift Sharing Rate", ">=25%", "of wishlist users", "Social commerce uptake"),
    ("Gift Fulfilment", ">=50%", "of shared items", "Collaborative shopping proof"),
    ("Re-engagement CTR", ">=12%", "open rate on alerts", "Price-drop notification effectiveness"),
    ("AOV Lift", "+$10", "vs. non-wishlist users", "Incremental spend proof"),
    ("NPS Impact", "+5 pts", "vs. non-wishlist users", "Customer satisfaction"),
]

for i, (name, target, unit, why) in enumerate(kpis):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 1.35)
    add_box(s, x, y, Inches(5.8), Inches(1.2), LIGHT_GREY)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(3.0), Inches(0.4), name, 15, CT_DARK, True)
    add_text(s, x + Inches(3.5), y + Inches(0.05), Inches(2.0), Inches(0.5), target, 24, CT_RED, True, PP_ALIGN.RIGHT)
    add_text(s, x + Inches(3.5), y + Inches(0.55), Inches(2.0), Inches(0.3), unit, 11, CT_GREY, alignment=PP_ALIGN.RIGHT)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(3.0), Inches(0.6), why, 11, CT_GREY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15: Recommendation
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Recommendation & Next Steps")

add_box(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.2), LIGHT_GREY)
add_text(s, Inches(1.0), Inches(1.5), Inches(11), Inches(1.0),
    "We recommend proceeding from POC to production build using EPAM EliteA, "
    "targeting pilot launch at 50 Canadian Tire locations within 13 weeks, "
    "with full rollout across all banners within 6 months.",
    16, CT_DARK, True, PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(2.9), Inches(11), Inches(0.5), "Next Steps", 22, CT_RED, True)

steps = [
    ("1", "Approve POC findings", "Stakeholder alignment on business case and revenue model", "Week 1", False),
    ("2", "Stand up EliteA team", "2 senior engineers + AI agent cluster, CTC API sandbox access", "Week 2", False),
    ("3", "Architecture & integration design", "API contracts with CTC backend, Triangle, product catalog", "Weeks 2-4", False),
    ("G", "ARCHITECTURE REVIEW GATE", "CTC & EPAM architecture board sign-off before development begins", "Week 4", True),
    ("4", "Agentic development sprint", "Parallel build: backend APIs + mobile app + test suites", "Weeks 5-10", False),
    ("5", "Integration & hardening", "CTC backend integration, security audit, performance testing", "Weeks 10-13", False),
    ("6", "Pilot launch (50 stores)", "Staged rollout with A/B testing, real-time monitoring", "Weeks 13-15", False),
]

for i, (num, title, desc, timeline, is_gate) in enumerate(steps):
    y = Inches(3.3 + i * 0.53)
    gate_color = RGBColor(0xF5, 0x7C, 0x00) if is_gate else CT_RED
    add_box(s, Inches(0.8), y, Inches(0.5), Inches(0.45), gate_color)
    add_text(s, Inches(0.8), y + Inches(0.02), Inches(0.5), Inches(0.45), num, 14, WHITE, True, PP_ALIGN.CENTER)
    title_color = RGBColor(0xF5, 0x7C, 0x00) if is_gate else CT_DARK
    add_text(s, Inches(1.5), y + Inches(0.02), Inches(3.5), Inches(0.42), title, 13 if is_gate else 14, title_color, True)
    add_text(s, Inches(5.0), y + Inches(0.02), Inches(5.0), Inches(0.42), desc, 11, CT_GREY)
    add_text(s, Inches(10.5), y + Inches(0.02), Inches(2.0), Inches(0.42), timeline, 11, gate_color, True, PP_ALIGN.RIGHT)

# Big numbers at bottom
add_box(s, Inches(0.8), Inches(6.7) - Inches(0.2), Inches(11.5), Inches(0.8), CT_DARK)
stats = [
    ("$189.5M", "CAD revenue uplift/yr"),
    ("$408K", "CAD EliteA investment"),
    ("15 weeks", "to pilot launch"),
    ("16,150%", "Year 1 ROI"),
]
for i, (val, label) in enumerate(stats):
    x = Inches(1.2 + i * 2.9)
    add_text(s, x, Inches(6.55), Inches(2.5), Inches(0.4), val, 20, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x, Inches(6.9), Inches(2.5), Inches(0.3), label, 11, RGBColor(0xCC,0xCC,0xCC), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# APPENDIX SLIDE A: Consumer Benefits Deep Dive
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Appendix A: Consumer Benefits — Deep Dive")

consumer_deep = [
    ("Save for Later — Reducing Decision Fatigue",
     "Shoppers frequently browse without immediate purchase intent. Today, if a CTC customer finds "
     "an interesting product in-store or online but isn't ready to buy, they have no way to save it "
     "within the CTC ecosystem. The item is lost. A wishlist creates a persistent, personal catalog "
     "of intent — items the customer wants but hasn't committed to yet. This reduces the cognitive "
     "load of remembering products, eliminates the friction of re-finding them later, and keeps the "
     "customer anchored to CTC rather than searching elsewhere. Industry data shows wishlists increase "
     "return visit rates by 30-50%, because the customer now has a reason to come back."),
    ("In-Store to Digital Bridge — Barcode Scanning",
     "Canadian Tire's 500+ stores are a massive discovery engine — customers touch, compare, and evaluate "
     "products in person. But that in-store discovery currently dead-ends. A customer who sees a Mastercraft "
     "drill on the shelf but wants to think about it has no way to capture that moment digitally. Barcode "
     "scanning bridges this gap: point the phone, scan the shelf tag, and the product is saved to a wishlist "
     "in seconds. The customer can then purchase online, at a different store, or during a future visit. "
     "This is especially powerful for high-consideration items (power tools, sporting equipment, appliances) "
     "where the purchase cycle spans days or weeks."),
    ("Gift Coordination — Eliminating Duplicate Gifts",
     "Gift-giving is a $30B+ annual category in Canadian retail, with major peaks at Christmas, birthdays, "
     "Father's/Mother's Day, and back-to-school. The #1 pain point for gift buyers is uncertainty — \"What do "
     "they actually want?\" and \"Has someone else already bought this?\" Shared wishlists solve both problems. "
     "The recipient curates exactly what they want from CTC's catalog. Multiple gift-givers can see the list "
     "and claim items, ensuring no duplicates. This drives full-price purchases (gift-givers rarely comparison-shop "
     "for specific wished items) and significantly higher fulfilment rates — industry data shows 60-70% of shared "
     "wishlist items get purchased vs. ~15% conversion on generic browsing."),
    ("Price Monitoring — Sale Alert Re-engagement",
     "Canadian Tire runs frequent promotional cycles — weekly flyers, seasonal sales, clearance events. A wishlist "
     "transforms these promotions from spray-and-pray to precision-targeted. When a wishlisted item goes on sale, "
     "the customer gets a personalized notification: \"Your Yardworks pressure washer is now 25% off.\" This is "
     "5-8x more effective than generic promotional emails because the intent already exists — the customer already "
     "told you they want this item. For CTC, this means higher promotion ROI, reduced promotional waste, and "
     "incremental conversions that wouldn't happen without the trigger."),
    ("Cross-Device Continuity — Seamless Shopping Journey",
     "Modern retail journeys span multiple devices and touchpoints. A customer might scan a barcode in-store on "
     "their phone, review the wishlist on a tablet at home, and complete the purchase on a desktop computer. "
     "Wishlist data synced across devices (via Triangle account) ensures no friction in this journey. The "
     "customer never has to re-find a product or remember which store they saw it in. This is particularly "
     "valuable for CTC's cross-banner ecosystem — a customer could add a Mark's jacket, a Canadian Tire drill, "
     "and SportChek running shoes to the same wishlist, then check out items from each banner seamlessly."),
]

for i, (title, detail) in enumerate(consumer_deep):
    y = Inches(1.3 + i * 1.2)
    add_text(s, Inches(0.8), y, Inches(12), Inches(0.35), title, 13, CT_RED, True)
    add_text(s, Inches(0.8), y + Inches(0.32), Inches(12), Inches(0.85), detail, 9, CT_DARK)

# ═══════════════════════════════════════════════════════════════════
# APPENDIX SLIDE B: CTC Business Benefits Deep Dive
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Appendix B: CTC Business Benefits — Deep Dive")

business_deep = [
    ("First-Party Intent Data — What Customers WANT",
     "Purchase data tells CTC what customers bought. Wishlist data tells CTC what customers want but haven't "
     "bought yet — a fundamentally different and more valuable signal. This intent data enables: (a) targeted "
     "promotions with known ROI (promote a product to people who already want it), (b) merchandising insights "
     "(which products are wishlisted most but purchased least — price barrier? availability?), (c) ad targeting "
     "for CTC's retail media network (brands pay premium to target high-intent audiences), and (d) demand "
     "forecasting with forward-looking signals rather than backward-looking sales data."),
    ("Triangle Rewards Enrichment — Canada's Richest Intent Dataset",
     "CTC's 16M Triangle members already generate purchase, loyalty, and demographic data. Layering wishlist "
     "intent data on top creates an unmatched customer intelligence platform. CTC can build predictive models: "
     "\"Customers who wishlisted X also purchased Y within 30 days.\" This powers personalized recommendations, "
     "dynamic pricing experiments, and a retail media offering that competes with Amazon's ad platform in the "
     "Canadian market. No other Canadian retailer has both the loyalty scale (16M) and the cross-category breadth "
     "(automotive, sporting goods, apparel, celebrations, outdoor) to build this."),
    ("Seasonal Revenue Capture — Converting Browsing to Pipelines",
     "CTC's revenue has significant seasonal peaks: Christmas (Nov-Dec), Spring/Summer (outdoor, garden), "
     "Back-to-School (Aug-Sep), and Father's Day. Wishlists convert casual seasonal browsing into committed "
     "purchase pipelines weeks before the actual buying moment. A parent browsing toys in October creates a "
     "Christmas wishlist that gets shared to grandparents in November, who buy in December. This extends the "
     "effective selling window, smooths demand, and gives CTC visibility into upcoming seasonal demand — "
     "invaluable for inventory planning at the dealer level."),
    ("Dealer-Owner Inventory Signals — Local Demand Intelligence",
     "CTC's unique dealer-owner model means each store makes its own stocking decisions. Aggregated wishlist "
     "data at the store/region level gives dealers a new signal: \"These products are being wishlisted in your "
     "area but you don't stock them.\" This reduces both overstock (carrying products nobody wants) and "
     "stockouts (missing products people are actively saving). For PartSource especially, where parts are "
     "high-SKU and location-dependent, wishlist data could significantly improve inventory efficiency."),
    ("Cross-Banner Discovery — Unified CTC Ecosystem",
     "A unified wishlist spanning Canadian Tire, SportChek, Mark's, Helly Hansen, Party City, and PartSource "
     "is something no competitor can replicate. A customer planning a camping trip could wishlist a Coleman tent "
     "(CT), hiking boots (SportChek), a Helly Hansen jacket, and Mark's base layers — all in one list. This "
     "drives cross-banner traffic that CTC currently has no mechanism to generate. Each banner benefits from "
     "discovery driven by the others, increasing total basket size across the CTC ecosystem."),
    ("Competitive Moat — First-Mover in Canadian Retail",
     "No Canadian general-merchandise retailer currently offers barcode-scan-to-wishlist combined with gift "
     "sharing and cross-banner support. Amazon has wishlists but no barcode scanning (no physical stores in "
     "Canada). Walmart has basic save-for-later but no sharing or barcode scan. Costco has nothing. By "
     "launching first, CTC sets the standard for wishlist commerce in Canadian retail, builds a data moat "
     "that grows with each user, and creates a feature that's deeply integrated with its physical store "
     "advantage — something pure e-commerce competitors cannot replicate."),
]

for i, (title, detail) in enumerate(business_deep):
    y = Inches(1.3 + i * 1.0)
    add_text(s, Inches(0.8), y, Inches(12), Inches(0.35), title, 13, SUCCESS_GREEN, True)
    add_text(s, Inches(0.8), y + Inches(0.3), Inches(12), Inches(0.7), detail, 9, CT_DARK)

# ═══════════════════════════════════════════════════════════════════
# APPENDIX SLIDE C: Success Metrics — Sources & Methodology
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Appendix C: Success Metrics — Sources & Methodology")

kpi_sources = [
    ("Wishlist Adoption >= 20% of MAU",
     "Based on Barilliance (2023) and Baymard Institute (2024) e-commerce feature adoption studies. "
     "Mature wishlist features in retail apps see 20-35% adoption among active users. We use the low end "
     "as a Year 1 target since the feature will be new. Amazon's wishlist adoption is estimated at ~30% of MAU; "
     "we discount for CTC's less digitally-native customer base."),
    ("Avg Items per Wishlist >= 5",
     "Industry benchmark from Monetate (2023) and Shopify merchant data (2024). Retail wishlists with active "
     "usage average 4-8 items. Lower for general merchandise (3-5), higher for specialty/gift registries (8-15). "
     "Target of 5 reflects a healthy engagement level without requiring power-user behaviour."),
    ("Wishlist-to-Purchase Conversion >= 15%",
     "Barilliance (2023): wishlist-to-cart conversion rates range from 12-25% across retail e-commerce. "
     "Monetate Digital Intelligence Report (2024): wishlisted items convert at 2-3x the rate of browsed-only items. "
     "15% is the midpoint for general merchandise retail."),
    ("Gift Sharing Rate >= 25% of wishlist users",
     "The Knot / WeddingWire registry data (2023): 40-60% of registry creators share with 5+ people. "
     "Amazon Wish List sharing: ~20-30% of list creators share at least once. We target 25% as a conservative "
     "estimate, reflecting that CTC's feature is new and sharing requires active intent."),
    ("Gift Fulfilment >= 50% of shared items",
     "National Retail Federation (2023): gift registries see 55-75% fulfilment rates for shared lists. "
     "Wedding registries: 65-85%. Birthday/holiday wishlists: 40-60%. Target of 50% reflects the mixed "
     "occasion types CTC would see (not just weddings, which skew higher)."),
    ("Re-engagement CTR >= 12%",
     "Braze (2024) and Iterable (2023) push notification benchmarks: generic retail push CTR is 2-4%. "
     "Personalized price-drop alerts achieve 10-18% CTR (Omnisend 2024). Sailthru (2023): triggered emails "
     "for wishlisted items see 5-8x baseline CTR. Target of 12% is mid-range for personalized alerts."),
    ("AOV Lift >= +$10 vs. non-wishlist users",
     "Barilliance (2023): wishlist users show 15-25% higher AOV vs. non-users. At CTC's estimated $55 "
     "average transaction value (blended across banners), a 15% lift = $8.25 and 25% = $13.75. "
     "Target of +$10 is the midpoint. Driven by higher purchase intent and larger basket building."),
    ("NPS Impact >= +5 points",
     "Qualtrics (2024) and Bain & Company: features that reduce friction and enable personalization "
     "typically improve NPS by 3-8 points. Target of +5 reflects the convenience value of wishlists "
     "(save for later, gift coordination) which directly addresses common shopper pain points."),
]

for i, (metric, source) in enumerate(kpi_sources):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 1.45)
    add_text(s, x, y, Inches(6.0), Inches(0.3), metric, 11, CT_RED, True)
    add_text(s, x, y + Inches(0.3), Inches(6.0), Inches(1.1), source, 8, CT_DARK)

# ═══════════════════════════════════════════════════════════════════
# Thank You
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_box(s, Inches(0), Inches(0), W, Inches(3.5), CT_RED)
add_text(s, Inches(0.8), Inches(1.0), Inches(11), Inches(1.2), "Thank You", 52, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(2.3), Inches(11), Inches(0.6),
    "Scan. Save. Share. — Let's build the future of Canadian Tire shopping together.",
    20, RGBColor(0xFF,0xCC,0xCC), alignment=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
    "EPAM Systems  |  Hackathon POC Team", 18, CT_DARK, True, PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
    "Powered by EPAM EliteA — Agentic AI Software Development", 16, CT_RED, False, PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
    "Questions?", 24, CT_GREY, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "CTC_Mobile_Wishlist_Business_Case.pptx")
prs.save(out_path)
print(f"\nSaved: {out_path}")
print(f"Slides: {len(prs.slides)}")
